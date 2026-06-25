"""Validate pipeline output."""

from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation

from langsmith import traceable

from src.extract.layout_extractor import extract_layout
from src.models import BlockValidation, ChartBlock, ContentMapping, LayoutDoc, ValidationReport


def _count_shapes(slide) -> int:
    count = 0

    def walk(shapes):
        nonlocal count
        for shape in shapes:
            count += 1
            if shape.shape_type.name == "GROUP":
                walk(shape.shapes)

    walk(slide.shapes)
    return count


def _chart_structure_ok(mapping, original: ChartBlock) -> bool:
    if len(mapping.categories) != original.category_count:
        return False
    if len(mapping.series) != original.series_count:
        return False
    return all(len(series.values) == original.category_count for series in mapping.series)


@traceable(name="validate_output", run_type="chain")
def validate_output(
    layout: LayoutDoc,
    mapping: ContentMapping,
    output_pptx: str | Path,
    input_pptx: str | Path | None = None,
) -> ValidationReport:
    output_pptx = Path(output_pptx)
    mapping_lookup = {m.block_id: m for m in mapping.mappings}
    chart_mapping_lookup = {m.block_id: m for m in mapping.chart_mappings}
    block_results: list[BlockValidation] = []
    unmapped: list[str] = []
    empty: list[str] = []
    warnings: list[str] = []
    manual_review: list[str] = []

    total_blocks = 0
    mapped_blocks = 0
    total_charts = 0
    mapped_charts = 0

    for slide in layout.slides:
        for block in slide.blocks:
            total_blocks += 1
            entry = mapping_lookup.get(block.block_id)
            mapped = entry is not None
            new_text = entry.new_text if entry else ""
            if mapped:
                mapped_blocks += 1
            else:
                unmapped.append(block.block_id)
            is_empty = mapped and not new_text.strip()
            if is_empty:
                empty.append(block.block_id)
            within_limit = len(new_text) <= block.char_limit if mapped else False
            if mapped and not within_limit:
                warnings.append(
                    f"{block.block_id}: new text length {len(new_text)} exceeds char_limit {block.char_limit}"
                )
            block_results.append(
                BlockValidation(
                    block_id=block.block_id,
                    original_len=len(block.original_text),
                    new_len=len(new_text) if mapped else 0,
                    within_char_limit=within_limit,
                    mapped=mapped,
                    empty=is_empty,
                    block_kind="text",
                )
            )

        for chart in slide.charts:
            total_charts += 1
            entry = chart_mapping_lookup.get(chart.block_id)
            mapped = entry is not None
            if mapped:
                mapped_charts += 1
                if not _chart_structure_ok(entry, chart):
                    warnings.append(
                        f"{chart.block_id}: chart mapping has wrong category/series counts"
                    )
            else:
                unmapped.append(chart.block_id)
            block_results.append(
                BlockValidation(
                    block_id=chart.block_id,
                    original_len=len(chart.title),
                    new_len=len(entry.title) if entry else 0,
                    within_char_limit=True,
                    mapped=mapped,
                    empty=mapped and not entry.title and not entry.categories,
                    block_kind="chart",
                )
            )

        for shape in slide.non_text_shapes:
            if shape.shape_type == "PICTURE":
                manual_review.append(
                    f"Slide {slide.slide_index}: PICTURE at {shape.block_id} requires manual review if content should change"
                )
            elif shape.shape_type == "CHART":
                manual_review.append(
                    f"Slide {slide.slide_index}: unsupported CHART at {shape.block_id} (non-category chart)"
                )

    shape_count_unchanged = True
    if input_pptx and Path(input_pptx).exists():
        try:
            in_prs = Presentation(str(input_pptx))
            out_prs = Presentation(str(output_pptx))
            if len(in_prs.slides) == len(out_prs.slides):
                for in_slide, out_slide in zip(in_prs.slides, out_prs.slides):
                    if _count_shapes(in_slide) != _count_shapes(out_slide):
                        shape_count_unchanged = False
                        break
            else:
                shape_count_unchanged = False
        except Exception:
            shape_count_unchanged = False

    editable_check = False
    try:
        with zipfile.ZipFile(output_pptx, "r") as zf:
            names = zf.namelist()
            editable_check = any(n.startswith("ppt/slides/slide") for n in names)
        out_layout = extract_layout(output_pptx)
        has_content = any(
            block.original_text for slide in out_layout.slides for block in slide.blocks
        ) or any(slide.charts for slide in out_layout.slides)
        editable_check = editable_check and has_content
    except Exception:
        editable_check = False

    total_items = total_blocks + total_charts
    mapped_items = mapped_blocks + mapped_charts
    coverage = (mapped_items / total_items * 100) if total_items else 100.0
    passed = (
        coverage >= 95.0
        and not empty
        and editable_check
        and shape_count_unchanged
    )

    return ValidationReport(
        passed=passed,
        coverage_pct=round(coverage, 2),
        mapped_blocks=mapped_blocks,
        total_text_blocks=total_blocks,
        mapped_charts=mapped_charts,
        total_charts=total_charts,
        unmapped_block_ids=unmapped,
        empty_block_ids=empty,
        block_results=block_results,
        shape_count_unchanged=shape_count_unchanged,
        editable_check=editable_check,
        manual_review_needed=manual_review,
        warnings=warnings,
    )
