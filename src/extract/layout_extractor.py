"""Recursive layout extractor for PPTX slides."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

from langsmith import traceable

from src.extract.chart_extractor import extract_chart_block
from src.font_color import _best_color_from_runs
from src.models import ChartBlock, LayoutDoc, NonTextShape, Role, SlideLayout, StyleHint, TextBlock

EMU_PER_INCH = 914400
DEFAULT_FONT_PT = 18.0


def _emu_to_inches(emu: int | Emu | None) -> float:
    if emu is None:
        return 0.0
    return int(emu) / EMU_PER_INCH


def _alignment_name(alignment) -> str | None:
    if alignment is None:
        return None
    mapping = {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
    }
    return mapping.get(alignment, str(alignment))


def _placeholder_type_name(shape) -> str | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        ph_type = shape.placeholder_format.type
        return ph_type.name if ph_type is not None else None
    except Exception:
        return None


def _shape_type_name(shape) -> str:
    try:
        return shape.shape_type.name if shape.shape_type is not None else "UNKNOWN"
    except Exception:
        return "UNKNOWN"


def _dominant_font_pt(shape) -> float:
    sizes: list[float] = []
    if not shape.has_text_frame:
        return DEFAULT_FONT_PT
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
    return max(sizes) if sizes else DEFAULT_FONT_PT


def _get_text_frame(obj):
    """Return text_frame from a shape or table cell."""
    if hasattr(obj, "has_text_frame"):
        return obj.text_frame if obj.has_text_frame else None
    return getattr(obj, "text_frame", None)


def _extract_style_hint(shape) -> StyleHint:
    text_frame = _get_text_frame(shape)
    if text_frame is None:
        return StyleHint()

    all_runs = [
        run
        for para in text_frame.paragraphs
        for run in para.runs
        if run.text.strip()
    ]
    if not all_runs:
        all_runs = [
            run for para in text_frame.paragraphs for run in para.runs
        ]

    color_rgb, color_theme, color_type, color_brightness = _best_color_from_runs(all_runs or [])
    dominant = all_runs[0] if all_runs else None
    first_para = text_frame.paragraphs[0] if text_frame.paragraphs else None

    if dominant:
        return StyleHint(
            font_name=dominant.font.name,
            size_pt=dominant.font.size.pt if dominant.font.size else DEFAULT_FONT_PT,
            bold=dominant.font.bold,
            italic=dominant.font.italic,
            alignment=_alignment_name(
                next((p.alignment for p in text_frame.paragraphs if p.runs), None)
                or (first_para.alignment if first_para else None)
            ),
            color_rgb=color_rgb,
            color_theme=color_theme,
            color_type=color_type,
            color_brightness=color_brightness,
        )

    return StyleHint(alignment=_alignment_name(first_para.alignment) if first_para else None)


def _extract_text_content(shape) -> tuple[str, int, list[int]]:
    if not shape.has_text_frame:
        return "", 0, []
    paragraphs = shape.text_frame.paragraphs
    lines: list[str] = []
    bullet_levels: list[int] = []
    for para in paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if text or len(paragraphs) == 1:
            lines.append(text)
            bullet_levels.append(para.level)
    text = "\n".join(lines).strip()
    return text, len(paragraphs) if paragraphs else 0, bullet_levels


def _estimate_char_limit(width_in: float, height_in: float, font_pt: float) -> int:
    if width_in <= 0 or height_in <= 0:
        return max(40, int(font_pt * 4))
    chars_per_line = max(8, int(width_in * 72 / (font_pt * 0.55)))
    lines = max(1, int(height_in * 72 / (font_pt * 1.25)))
    return max(20, chars_per_line * lines)


def _role_from_placeholder(ph_type: str | None) -> Role | None:
    if not ph_type:
        return None
    upper = ph_type.upper()
    if "TITLE" in upper and "SUBTITLE" not in upper:
        return "title"
    if "SUBTITLE" in upper:
        return "subtitle"
    if "BODY" in upper or "CONTENT" in upper or "OBJECT" in upper:
        return "body"
    return None


def _infer_role(
    shape,
    text: str,
    bullet_levels: list[int],
    font_pt: float,
    top_in: float,
    placeholder_type: str | None,
    is_table_cell: bool = False,
) -> Role:
    if is_table_cell:
        return "table_cell"
    ph_role = _role_from_placeholder(placeholder_type)
    if ph_role:
        return ph_role
    if bullet_levels and any(level > 0 for level in bullet_levels):
        return "body"
    if bullet_levels and len(bullet_levels) > 1:
        return "body"
    if top_in < 1.5 and font_pt >= 24 and len(text) < 80:
        return "title"
    if top_in < 2.0 and font_pt >= 20 and len(text) < 100:
        return "subtitle"
    if len(text) <= 30 and font_pt <= 16:
        return "label"
    if bullet_levels:
        return "body"
    return "other"


class _SlideWalker:
    def __init__(self, slide_index: int):
        self.slide_index = slide_index
        self.blocks: list[TextBlock] = []
        self.charts: list[ChartBlock] = []
        self.non_text_shapes: list[NonTextShape] = []
        self._text_candidates: list[tuple[float, float, TextBlock]] = []

    def walk(self, shapes, path: str = "shapes") -> None:
        for idx, shape in enumerate(shapes):
            current_path = f"{path}[{idx}]"
            block_id = f"{self.slide_index}/{current_path}"

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.walk(shape.shapes, f"{current_path}/group")
                continue

            geom = {
                "left_in": _emu_to_inches(shape.left),
                "top_in": _emu_to_inches(shape.top),
                "width_in": _emu_to_inches(shape.width),
                "height_in": _emu_to_inches(shape.height),
            }

            if shape.has_table:
                self._extract_table(shape, current_path, geom)
                continue

            if shape.has_text_frame:
                text, para_count, bullet_levels = _extract_text_content(shape)
                if text:
                    font_pt = _dominant_font_pt(shape)
                    placeholder_type = _placeholder_type_name(shape)
                    role = _infer_role(
                        shape,
                        text,
                        bullet_levels,
                        font_pt,
                        geom["top_in"],
                        placeholder_type,
                    )
                    block = TextBlock(
                        block_id=block_id,
                        role=role,
                        original_text=text,
                        char_limit=_estimate_char_limit(
                            geom["width_in"], geom["height_in"], font_pt
                        ),
                        paragraph_count=max(1, para_count),
                        bullet_levels=bullet_levels or [0],
                        style_hint=_extract_style_hint(shape),
                        is_placeholder=getattr(shape, "is_placeholder", False),
                        placeholder_type=placeholder_type,
                        **geom,
                    )
                    self.blocks.append(block)
                    self._text_candidates.append((geom["top_in"], -font_pt, block))
                continue

            if shape.has_chart:
                chart_block = extract_chart_block(shape, block_id, geom)
                if chart_block:
                    self.charts.append(chart_block)
                else:
                    self.non_text_shapes.append(
                        NonTextShape(
                            block_id=block_id,
                            shape_type="CHART",
                            editable=False,
                            **geom,
                        )
                    )
                continue

            shape_type = _shape_type_name(shape)
            if shape_type in {"PICTURE", "LINE", "FREEFORM", "AUTO_SHAPE"}:
                editable = shape_type not in {"PICTURE", "LINE"}
                self.non_text_shapes.append(
                    NonTextShape(
                        block_id=block_id,
                        shape_type=shape_type,
                        editable=editable,
                        **geom,
                    )
                )

    def _extract_table(self, shape, path: str, geom: dict) -> None:
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_path = f"{path}/table[{row_idx},{col_idx}]"
                block_id = f"{self.slide_index}/{cell_path}"
                text = cell.text.strip()
                if not text:
                    continue
                font_pt = DEFAULT_FONT_PT
                if cell.text_frame.paragraphs:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                font_pt = run.font.size.pt
                                break
                cell_geom = dict(geom)
                block = TextBlock(
                    block_id=block_id,
                    role="table_cell",
                    original_text=text,
                    char_limit=_estimate_char_limit(
                        geom["width_in"] / max(1, len(table.columns)),
                        geom["height_in"] / max(1, len(table.rows)),
                        font_pt,
                    ),
                    paragraph_count=len(cell.text_frame.paragraphs) or 1,
                    bullet_levels=[0],
                    style_hint=_extract_style_hint(cell),
                    **cell_geom,
                )
                self.blocks.append(block)

    def refine_title_roles(self) -> None:
        """Promote top-most large text block to title if none assigned."""
        if any(b.role == "title" for b in self.blocks):
            return
        if not self._text_candidates:
            return
        self._text_candidates.sort(key=lambda item: (item[0], item[1]))
        top_block = self._text_candidates[0][2]
        for block in self.blocks:
            if block.block_id == top_block.block_id and block.role == "other":
                block.role = "title"


@traceable(name="extract_layout", run_type="chain")
def extract_layout(input_path: str | Path) -> LayoutDoc:
    input_path = Path(input_path)
    prs = Presentation(str(input_path))
    slides: list[SlideLayout] = []

    for slide_index, slide in enumerate(prs.slides):
        walker = _SlideWalker(slide_index)
        walker.walk(slide.shapes)
        walker.refine_title_roles()
        slides.append(
            SlideLayout(
                slide_index=slide_index,
                layout_name=slide.slide_layout.name,
                slide_xml_path=f"ppt/slides/slide{slide_index + 1}.xml",
                blocks=walker.blocks,
                charts=walker.charts,
                non_text_shapes=walker.non_text_shapes,
            )
        )

    return LayoutDoc(source_file=str(input_path), slides=slides)
