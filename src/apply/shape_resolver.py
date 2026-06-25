"""Resolve block_id paths to shapes on a slide."""

from __future__ import annotations

import re

from pptx.enum.shapes import MSO_SHAPE_TYPE


def resolve_shape(slide, path: str):
    """Resolve a path like shapes[0]/group[1] to the target shape."""
    segments = path.split("/")
    current_shapes = slide.shapes
    shape = None

    for segment in segments:
        shape_match = re.match(r"^shapes\[(\d+)\]$", segment)
        if shape_match:
            shape = current_shapes[int(shape_match.group(1))]
            continue

        group_match = re.match(r"^group\[(\d+)\]$", segment)
        if group_match:
            if shape is None or shape.shape_type != MSO_SHAPE_TYPE.GROUP:
                raise ValueError(f"Expected group shape for segment {segment}")
            current_shapes = shape.shapes
            shape = current_shapes[int(group_match.group(1))]
            continue

        table_match = re.match(r"^table\[(\d+),(\d+)\]$", segment)
        if table_match:
            if shape is None or not shape.has_table:
                raise ValueError(f"Expected table shape for segment {segment}")
            row, col = int(table_match.group(1)), int(table_match.group(2))
            return shape.table.cell(row, col)

        raise ValueError(f"Unknown path segment: {segment}")

    if shape is None:
        raise ValueError(f"Could not resolve shape at path {path}")
    return shape
