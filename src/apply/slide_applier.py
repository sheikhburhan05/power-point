"""Apply content mapping back to PPTX template."""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from langsmith import traceable

from src.apply.chart_applier import apply_chart_mapping
from src.apply.text_replacer import apply_autofit, replace_text_in_frame
from src.models import ChartBlock, ContentMapping, LayoutDoc, TextBlock


def _parse_block_id(block_id: str) -> tuple[int, str]:
    match = re.match(r"^(\d+)/(.*)$", block_id)
    if not match:
        raise ValueError(f"Invalid block_id: {block_id}")
    return int(match.group(1)), match.group(2)


def _resolve_text_frame(slide, path: str):
    segments = path.split("/")
    current_shapes = slide.shapes
    shape = None

    for segment in segments:
        shape_match = re.match(r"^shapes\[(\d+)\]$", segment)
        if shape_match:
            shape = current_shapes[int(shape_match.group(1))]
            continue

        group_match = re.match(r"^group\[(\d+)\]$", segment)
        if group_match:
            if shape is None or shape.shape_type != MSO_SHAPE_TYPE.GROUP:
                raise ValueError(f"Expected group shape for segment {segment}")
            current_shapes = shape.shapes
            shape = current_shapes[int(group_match.group(1))]
            continue

        table_match = re.match(r"^table\[(\d+),(\d+)\]$", segment)
        if table_match:
            if shape is None or not shape.has_table:
                raise ValueError(f"Expected table shape for segment {segment}")
            row, col = int(table_match.group(1)), int(table_match.group(2))
            return shape.table.cell(row, col).text_frame

        raise ValueError(f"Unknown path segment: {segment}")

    if shape is None or not shape.has_text_frame:
        raise ValueError(f"No text frame at path {path}")
    return shape.text_frame


def _block_lookup(layout: LayoutDoc) -> dict[str, TextBlock]:
    return {
        block.block_id: block
        for slide in layout.slides
        for block in slide.blocks
    }


def _chart_lookup(layout: LayoutDoc) -> dict[str, ChartBlock]:
    return {
        chart.block_id: chart
        for slide in layout.slides
        for chart in slide.charts
    }


@traceable(name="apply_content", run_type="chain")
def apply_content(
    input_path: str | Path,
    layout: LayoutDoc,
    mapping: ContentMapping,
    output_path: str | Path,
) -> Path:
    input_path = Path(input_path)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(input_path))
    block_lookup = _block_lookup(layout)
    chart_lookup = _chart_lookup(layout)

    for entry in mapping.mappings:
        block = block_lookup.get(entry.block_id)
        if not block:
            continue
        slide_idx, rel_path = _parse_block_id(entry.block_id)
        slide = prs.slides[slide_idx]
        try:
            text_frame = _resolve_text_frame(slide, rel_path)
            replace_text_in_frame(text_frame, entry.new_text, block)
            apply_autofit(text_frame, block.role)
        except Exception:
            continue

    for chart_mapping in mapping.chart_mappings:
        original = chart_lookup.get(chart_mapping.block_id)
        if not original:
            continue
        slide_idx, rel_path = _parse_block_id(chart_mapping.block_id)
        slide = prs.slides[slide_idx]
        try:
            apply_chart_mapping(slide, chart_mapping, rel_path)
        except Exception:
            continue

    prs.save(str(output_path))
    return output_path
