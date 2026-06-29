"""Run-level text replacement preserving style."""

from __future__ import annotations

from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt

from src.font_color import apply_font_color, extract_font_color
from src.models import StyleHint, TextBlock


def _apply_alignment(paragraph, alignment: str | None) -> None:
    if not alignment:
        return
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    if alignment in mapping:
        paragraph.alignment = mapping[alignment]


def _copy_font(source_font, target_font) -> None:
    if source_font.name:
        target_font.name = source_font.name
    if source_font.size:
        target_font.size = source_font.size
    if source_font.bold is not None:
        target_font.bold = source_font.bold
    if source_font.italic is not None:
        target_font.italic = source_font.italic
    color_rgb, color_theme, color_type, color_brightness = extract_font_color(source_font)
    apply_font_color(
        target_font, color_rgb, color_theme, color_type, color_brightness
    )


def _apply_style_to_run(run, style: StyleHint, fallback_run=None) -> None:
    if fallback_run and not style.color_type:
        _copy_font(fallback_run.font, run.font)
    if style.font_name:
        run.font.name = style.font_name
    if style.size_pt:
        run.font.size = Pt(style.size_pt)
    if style.bold is not None:
        run.font.bold = style.bold
    if style.italic is not None:
        run.font.italic = style.italic
    if style.color_type:
        apply_font_color(
            run.font,
            style.color_rgb,
            style.color_theme,
            style.color_type,
            style.color_brightness,
        )


def _clear_text_frame(text_frame) -> None:
    text_frame.clear()
    if not text_frame.paragraphs:
        text_frame.add_paragraph()


def _snapshot_run_styles(text_frame) -> list[list[StyleHint]]:
    """Capture per-paragraph run styles before clearing (for multi-color text)."""
    snapshot: list[list[StyleHint]] = []
    for para in text_frame.paragraphs:
        para_styles: list[StyleHint] = []
        for run in para.runs:
            color_rgb, color_theme, color_type, color_brightness = extract_font_color(run.font)
            para_styles.append(
                StyleHint(
                    font_name=run.font.name,
                    size_pt=run.font.size.pt if run.font.size else None,
                    bold=run.font.bold,
                    italic=run.font.italic,
                    color_rgb=color_rgb,
                    color_theme=color_theme,
                    color_type=color_type,
                    color_brightness=color_brightness,
                )
            )
        snapshot.append(para_styles)
    return snapshot


def replace_text_in_frame(text_frame, new_text: str, block: TextBlock) -> None:
    """Replace all text in a text frame while preserving styling."""
    paragraphs = new_text.split("\n")
    if not paragraphs:
        paragraphs = [""]

    expected = max(1, block.paragraph_count)
    while len(paragraphs) < expected:
        paragraphs.append("")
    if len(paragraphs) > expected and expected == 1:
        paragraphs = [" ".join(paragraphs)]

    run_style_snapshot = _snapshot_run_styles(text_frame)
    fallback_run = None
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        fallback_run = text_frame.paragraphs[0].runs[0]

    _clear_text_frame(text_frame)

    for idx, para_text in enumerate(paragraphs[: max(len(paragraphs), expected)]):
        if idx == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()
        level = block.bullet_levels[idx] if idx < len(block.bullet_levels) else 0
        para.level = level
        _apply_alignment(para, block.style_hint.alignment)

        old_para_styles = (
            run_style_snapshot[idx] if idx < len(run_style_snapshot) else []
        )

        if len(old_para_styles) > 1 and para_text:
            # Preserve multi-run colors by splitting text across original run count
            chunks = _split_text_for_runs(para_text, len(old_para_styles))
            for chunk, run_style in zip(chunks, old_para_styles):
                run = para.add_run()
                run.text = chunk
                merged = block.style_hint.model_copy(update={
                    "color_rgb": run_style.color_rgb or block.style_hint.color_rgb,
                    "color_theme": run_style.color_theme or block.style_hint.color_theme,
                    "color_type": run_style.color_type or block.style_hint.color_type,
                    "color_brightness": (
                        run_style.color_brightness
                        if run_style.color_brightness is not None
                        else block.style_hint.color_brightness
                    ),
                    "font_name": run_style.font_name or block.style_hint.font_name,
                    "size_pt": run_style.size_pt or block.style_hint.size_pt,
                    "bold": run_style.bold if run_style.bold is not None else block.style_hint.bold,
                    "italic": run_style.italic if run_style.italic is not None else block.style_hint.italic,
                })
                _apply_style_to_run(run, merged, fallback_run)
        else:
            run = para.add_run()
            run.text = para_text
            style = block.style_hint
            if old_para_styles:
                rs = old_para_styles[0]
                style = block.style_hint.model_copy(update={
                    "color_rgb": rs.color_rgb or block.style_hint.color_rgb,
                    "color_theme": rs.color_theme or block.style_hint.color_theme,
                    "color_type": rs.color_type or block.style_hint.color_type,
                    "color_brightness": (
                        rs.color_brightness
                        if rs.color_brightness is not None
                        else block.style_hint.color_brightness
                    ),
                })
            _apply_style_to_run(run, style, fallback_run)


def _split_text_for_runs(text: str, run_count: int) -> list[str]:
    """Split replacement text into run_count chunks proportionally."""
    if run_count <= 1:
        return [text]
    words = text.split()
    if not words:
        return [""] * run_count
    chunks: list[str] = []
    per_run = max(1, len(words) // run_count)
    start = 0
    for i in range(run_count):
        if i == run_count - 1:
            chunk_words = words[start:]
        else:
            chunk_words = words[start : start + per_run]
            start += per_run
        chunks.append(" ".join(chunk_words))
    return chunks


def apply_autofit(text_frame, role: str) -> None:
    text_frame.word_wrap = True
    if role in {"title", "label", "table_cell", "subtitle"}:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    elif role == "body":
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    else:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
