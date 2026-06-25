"""Create a sample input slide for demo/testing."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def create_sample_slide(output_path: str | Path) -> Path:
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Q2 Business Review"
    for run in title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True

    # Use body placeholder if available
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type.name in {"BODY", "OBJECT"}:
            body = shape
            break
    if body is None:
        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))

    tf = body.text_frame
    tf.clear()
    bullets = [
        "Revenue grew 12% year over year",
        "Operating costs reduced by 5%",
        "Expanded into 3 new regional markets",
    ]
    for i, text in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        para.level = 0
        for run in para.runs:
            run.font.size = Pt(18)

    prs.save(str(output_path))
    return output_path


if __name__ == "__main__":
    create_sample_slide("input/slide.pptx")
    print("Created input/slide.pptx")
