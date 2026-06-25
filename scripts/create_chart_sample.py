"""Create a sample slide with a chart for testing."""

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


def create_chart_sample(output_path: str | Path) -> Path:
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Q2 Business Review"
    for run in title_tf.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("Revenue", (10.5, 12.3, 11.8, 14.2))
    chart_data.add_series("Costs", (6.2, 6.0, 5.8, 6.1))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(1.5),
        Inches(8.4),
        Inches(4.5),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Quarterly Performance"

    prs.save(str(output_path))
    return output_path


if __name__ == "__main__":
    create_chart_sample("input/chart_sample.pptx")
    print("Created input/chart_sample.pptx")
